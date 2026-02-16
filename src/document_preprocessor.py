#!/usr/bin/env python3
"""
Document Preprocessor Module
Handles multiple document formats: VTT, DOCX, PPTX, PDF, TXT, MD
"""

import re
import os
from pathlib import Path
from typing import Dict, List, Optional
from dataclasses import dataclass, field

# Document processing libraries
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False


@dataclass
class ProcessedDocument:
    """Container for processed document data"""
    content: str
    metadata: Dict
    document_type: str
    speakers: List[str] = field(default_factory=list)
    word_count: int = 0
    sections: List[Dict] = field(default_factory=list)
    images: List[str] = field(default_factory=list)


class DocumentPreprocessor:
    """Process multiple document formats for newsletter generation"""
    
    def __init__(self):
        self.supported_formats = {
            '.vtt': self._process_vtt,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.pdf': self._process_pdf,
            '.txt': self._process_text,
            '.md': self._process_markdown
        }
    
    def process_document(self, file_path: str) -> ProcessedDocument:
        """Process document based on file extension"""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        extension = path.suffix.lower()
        
        if extension not in self.supported_formats:
            raise ValueError(f"Unsupported format: {extension}. Supported: {list(self.supported_formats.keys())}")
        
        print(f"📄 Processing {extension} file: {path.name}")
        processor = self.supported_formats[extension]
        return processor(file_path)
    
    def _process_vtt(self, file_path: str) -> ProcessedDocument:
        """Process VTT transcript files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        content_lines = []
        speakers = set()
        sections = []
        current_section = {"speaker": None, "content": [], "timestamp": None}
        
        for line in lines:
            line = line.strip()
            
            # Extract timestamp
            timestamp_match = re.match(r'(\d{2}:\d{2}:\d{2}\.\d{3}) --> (\d{2}:\d{2}:\d{2}\.\d{3})', line)
            if timestamp_match:
                current_section["timestamp"] = timestamp_match.group(1)
                continue
            
            # Extract speaker and content
            speaker_match = re.search(r'<v ([^>]+)>(.+?)</v>', line)
            if speaker_match:
                speaker = speaker_match.group(1).strip()
                content = speaker_match.group(2).strip()
                
                speakers.add(speaker)
                content_lines.append(f"{speaker}: {content}")
                
                if current_section["speaker"] != speaker and current_section["content"]:
                    sections.append(current_section.copy())
                    current_section = {"speaker": speaker, "content": [content], "timestamp": current_section["timestamp"]}
                else:
                    current_section["speaker"] = speaker
                    current_section["content"].append(content)
        
        if current_section["content"]:
            sections.append(current_section)
        
        full_content = "\n".join(content_lines)
        
        return ProcessedDocument(
            content=full_content,
            metadata={"file_name": Path(file_path).name, "format": "vtt", "source_type": "transcript"},
            document_type="transcript",
            speakers=sorted(list(speakers)),
            word_count=len(full_content.split()),
            sections=sections,
            images=[]
        )
    
    def _process_docx(self, file_path: str) -> ProcessedDocument:
        """Process Word documents"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
        
        doc = Document(file_path)
        
        content_lines = []
        sections = []
        current_section = {"heading": "Introduction", "content": []}
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Check if it's a heading
            if para.style.name.startswith('Heading'):
                if current_section["content"]:
                    sections.append(current_section.copy())
                current_section = {"heading": text, "content": []}
            else:
                current_section["content"].append(text)
                content_lines.append(text)
        
        if current_section["content"]:
            sections.append(current_section)
        
        full_content = "\n\n".join(content_lines)
        
        return ProcessedDocument(
            content=full_content,
            metadata={"file_name": Path(file_path).name, "format": "docx", "source_type": "document"},
            document_type="document",
            speakers=[],
            word_count=len(full_content.split()),
            sections=sections,
            images=[]
        )
    
    def _process_pptx(self, file_path: str) -> ProcessedDocument:
        """Process PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        prs = Presentation(file_path)
        
        content_lines = []
        sections = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_title = f"Slide {slide_num}"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_content.append(text)
                    
                    # First text box is usually the title
                    if len(slide_content) == 1:
                        slide_title = text
            
            if slide_content:
                sections.append({
                    "slide_number": slide_num,
                    "title": slide_title,
                    "content": slide_content
                })
                content_lines.extend(slide_content)
        
        full_content = "\n\n".join(content_lines)
        
        return ProcessedDocument(
            content=full_content,
            metadata={"file_name": Path(file_path).name, "format": "pptx", "source_type": "presentation", "slide_count": len(prs.slides)},
            document_type="presentation",
            speakers=[],
            word_count=len(full_content.split()),
            sections=sections,
            images=[]
        )
    
    def _process_pdf(self, file_path: str) -> ProcessedDocument:
        """Process PDF documents"""
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2 not installed. Install with: pip install PyPDF2")
        
        content_lines = []
        
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    content_lines.append(text.strip())
        
        full_content = "\n\n".join(content_lines)
        
        return ProcessedDocument(
            content=full_content,
            metadata={"file_name": Path(file_path).name, "format": "pdf", "source_type": "document"},
            document_type="document",
            speakers=[],
            word_count=len(full_content.split()),
            sections=[],
            images=[]
        )
    
    def _process_text(self, file_path: str) -> ProcessedDocument:
        """Process plain text files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return ProcessedDocument(
            content=content,
            metadata={"file_name": Path(file_path).name, "format": "txt", "source_type": "text"},
            document_type="text",
            speakers=[],
            word_count=len(content.split()),
            sections=[],
            images=[]
        )
    
    def _process_markdown(self, file_path: str) -> ProcessedDocument:
        """Process Markdown files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Extract sections based on headers
        sections = []
        current_section = {"heading": "Introduction", "content": []}
        
        for line in content.split('\n'):
            if line.startswith('#'):
                if current_section["content"]:
                    sections.append(current_section.copy())
                heading = re.sub(r'^#+\s*', '', line).strip()
                current_section = {"heading": heading, "content": []}
            else:
                if line.strip():
                    current_section["content"].append(line)
        
        if current_section["content"]:
            sections.append(current_section)
        
        return ProcessedDocument(
            content=content,
            metadata={"file_name": Path(file_path).name, "format": "md", "source_type": "markdown"},
            document_type="markdown",
            speakers=[],
            word_count=len(content.split()),
            sections=sections,
            images=[]
        )
    
    def combine_documents(self, docs: List[ProcessedDocument]) -> ProcessedDocument:
        """Combine multiple processed documents into one"""
        combined_content = []
        all_speakers = set()
        all_sections = []
        all_images = []
        total_words = 0
        
        for doc in docs:
            combined_content.append(f"\n\n=== Source: {doc.metadata['file_name']} ===\n\n")
            combined_content.append(doc.content)
            all_speakers.update(doc.speakers)
            all_sections.extend(doc.sections)
            all_images.extend(doc.images)
            total_words += doc.word_count
        
        return ProcessedDocument(
            content="\n".join(combined_content),
            metadata={"combined": True, "source_count": len(docs)},
            document_type="combined",
            speakers=sorted(list(all_speakers)),
            word_count=total_words,
            sections=all_sections,
            images=all_images
        )


if __name__ == "__main__":
    # Test the preprocessor
    preprocessor = DocumentPreprocessor()
    print("Document Preprocessor initialized successfully!")
    print(f"Supported formats: {list(preprocessor.supported_formats.keys())}")

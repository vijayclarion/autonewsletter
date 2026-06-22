"""PowerPoint (.pptx) parser: one chunk per slide (title + body + notes), slide provenance."""

from __future__ import annotations

from pptx import Presentation

from newsletter_engine.ingestion.chunker import ParsedChunk


def parse(path) -> list[ParsedChunk]:
    presentation = Presentation(str(path))
    chunks: list[ParsedChunk] = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notes: {notes}")
        if parts:
            chunks.append(ParsedChunk(text="\n".join(parts), location={"slide": number}))
    return chunks

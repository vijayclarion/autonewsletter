"""Generate the synthetic dry-run sample content in fixtures/ (T005).

Creates a Teams-style transcript .docx, a .pptx deck, a .pdf doc, a .txt note, a generic
non-transcript .docx, and a corrupt file. Run once: python scripts/make_fixtures.py
"""

from __future__ import annotations

from pathlib import Path

import docx
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
FIXTURES = ROOT / "fixtures"

TRANSCRIPT_TURNS = [
    ("Priya Sharma", "0:05", "Good morning everyone, how was the long weekend? I finally "
     "got around to that hiking trip with the kids."),
    ("Daniel Brooks", "0:21", "It was great, thanks! We did a barbecue on Sunday. Anyway, "
     "shall we get started?"),
    ("Priya Sharma", "0:48", "Yes. Today's main topic is the API gateway migration. We are "
     "moving from the legacy reverse proxy to a managed API gateway so we can centralize "
     "authentication, rate limiting, and request routing across all microservices."),
    ("Daniel Brooks", "1:32", "The proof of concept showed the gateway adds about three "
     "milliseconds of latency per request, which is acceptable. The bigger win is that "
     "JWT token validation moves out of each service into the gateway layer."),
    ("Aisha Khan", "2:10", "From the data side, we also validated the new caching strategy. "
     "Redis with a write-through cache cut median read latency from 120 milliseconds to 18, "
     "and the cache hit rate is holding at 94 percent in staging."),
    ("Daniel Brooks", "3:05", "One open decision: do we deploy the gateway per region or as "
     "a single global cluster? Per-region deployment doubles the infrastructure cost but "
     "removes the cross-region failover risk we saw during the March incident."),
    ("Priya Sharma", "3:50", "Let's take per-region as the working assumption. Action item "
     "for the platform team: produce a cost comparison of both topologies by next sprint."),
    ("Aisha Khan", "4:30", "By the way, is anyone going to the office party on Friday? I "
     "heard the new cafeteria is doing the catering."),
    ("Priya Sharma", "4:41", "I'll be there! Okay, back to work — Daniel, can you also "
     "schedule the security review of the gateway's OAuth integration with the identity "
     "platform team?"),
]

DECK_SLIDES = [
    ("Platform Architecture Review — May 2026",
     "Engineering monthly sync\nPlatform & Data teams"),
    ("Current State: Legacy Reverse Proxy",
     "Authentication duplicated in every microservice\n"
     "No centralized rate limiting\nManual TLS certificate rotation\n"
     "Routing rules spread across twelve config files"),
    ("Target State: Managed API Gateway",
     "Central JWT validation at the edge\nDeclarative routing with per-service policies\n"
     "Built-in rate limiting and request quotas\nAutomatic certificate management"),
    ("Caching Layer Results",
     "Redis write-through cache in staging\nMedian read latency: 120 ms to 18 ms\n"
     "Cache hit rate: 94 percent\nNext: production rollout behind a feature flag"),
]

PDF_TEXT_LINES = [
    "Design Note: Caching Strategy for the Read Path",
    "",
    "We evaluated three caching approaches for the product catalog read path:",
    "cache-aside, write-through, and write-behind. Write-through was selected",
    "because it keeps the cache consistent with the database on every write,",
    "which matters for pricing data. The Redis cluster runs three nodes with",
    "replication enabled. Invalidation uses a 15-minute TTL as a safety net.",
    "Estimated infrastructure cost is 220 USD per month in the staging tier.",
]

NOTES_TXT = """Deployment checklist notes — platform team

The gateway rollout needs a canary stage: route five percent of traffic through the new
API gateway for 48 hours, watch the p99 latency dashboard, then ramp to fifty percent.

Rollback plan: DNS swap back to the legacy proxy, which stays warm for two weeks.

Reminder: book the team lunch for the quarter celebration, somewhere with outdoor seating.
"""

TEAM_UPDATE_PARAGRAPHS = [
    ("Team Update — May", "Heading 1"),
    ("Welcome and farewells", "Heading 2"),
    ("This month we welcomed two new engineers to the platform team. We also said goodbye "
     "to our office manager, who is relocating. The team organized a farewell picnic at "
     "the lakeside park and everyone had a wonderful time.", None),
    ("Social corner", "Heading 2"),
    ("The football group now meets on Thursdays. The book club picked a new novel for June "
     "and the baking competition was won by the QA team with a spectacular lemon cake.", None),
    ("One engineering note", "Heading 2"),
    ("The build pipeline upgrade to the new CI runners finished this month, cutting average "
     "build time from 14 minutes to 6.", None),
]


def make_transcript(path: Path) -> None:
    document = docx.Document()
    document.add_paragraph("Architecture Sync — May 2026")
    document.add_paragraph("Meeting transcript")
    for speaker, ts, text in TRANSCRIPT_TURNS:
        document.add_paragraph(f"{speaker}   {ts}")
        document.add_paragraph(text)
    document.save(str(path))


def make_deck(path: Path) -> None:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]  # title + content
    for title, body in DECK_SLIDES:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    notes = presentation.slides[3].notes_slide
    notes.notes_text_frame.text = (
        "Mention that the feature flag rollout starts the first week of June."
    )
    presentation.save(str(path))


def make_pdf(path: Path) -> None:
    """Hand-assembled single-page PDF with a correct xref table."""
    def text_stream() -> bytes:
        lines = ["BT /F1 11 Tf 50 760 Td 14 TL"]
        for line in PDF_TEXT_LINES:
            escaped = line.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
            lines.append(f"({escaped}) Tj T*")
        lines.append("ET")
        return "\n".join(lines).encode("latin-1")

    stream = text_stream()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_pos = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n"
    ).encode()
    path.write_bytes(bytes(out))


def make_team_update(path: Path) -> None:
    document = docx.Document()
    for text, style in TEAM_UPDATE_PARAGRAPHS:
        if style:
            document.add_paragraph(text, style=style)
        else:
            document.add_paragraph(text)
    document.save(str(path))


def main() -> None:
    sync_dir = FIXTURES / "2026-05-architecture-sync"
    hr_dir = FIXTURES / "2026-05-team-update"
    sync_dir.mkdir(parents=True, exist_ok=True)
    hr_dir.mkdir(parents=True, exist_ok=True)

    make_transcript(sync_dir / "transcript.docx")
    make_deck(sync_dir / "platform-deck.pptx")
    make_pdf(sync_dir / "caching-design-note.pdf")
    (sync_dir / "deployment-notes.txt").write_text(NOTES_TXT, encoding="utf-8")
    make_team_update(hr_dir / "team-update.docx")
    (FIXTURES / "corrupt.docx").write_bytes(b"This is not a real docx file \x00\x01\x02garbage")
    print(f"Fixtures written to {FIXTURES}")


if __name__ == "__main__":
    main()
